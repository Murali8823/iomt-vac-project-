from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(25, 118, 210)  # Blue
ACCENT_COLOR = RGBColor(66, 165, 245)   # Light Blue
TEXT_COLOR = RGBColor(33, 33, 33)       # Dark Gray
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 248, 255)      # Alice Blue

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = PRIMARY_COLOR
    header_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = PRIMARY_COLOR
    header_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.4))
    left_header_frame = left_header.text_frame
    left_header_p = left_header_frame.paragraphs[0]
    left_header_p.text = left_title
    left_header_p.font.size = Pt(22)
    left_header_p.font.bold = True
    left_header_p.font.color.rgb = PRIMARY_COLOR
    
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.2), Inches(5))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(0.4))
    right_header_frame = right_header.text_frame
    right_header_p = right_header_frame.paragraphs[0]
    right_header_p.text = right_title
    right_header_p.font.size = Pt(22)
    right_header_p.font.bold = True
    right_header_p.font.color.rgb = PRIMARY_COLOR
    
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.2), Inches(5))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "🩺 VitalGuard AI", "Real-Time IoMT Health Monitoring System")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "• VitalGuard AI is a state-of-the-art multi-model AI framework",
    "• Designed for real-time monitoring and anomaly detection in Medical IoT (IoMT)",
    "• Fuses time-series forecasting, edge-based anomaly detection, and deep learning",
    "• Ensures patient safety through continuous, intelligent oversight",
    "• Built with FastAPI backend and modern web dashboard frontend"
])

# Slide 3: Key Features
add_two_column_slide(prs, "Key Features",
    "Core Capabilities",
    [
        "📡 Real-Time Simulation: Generates vital sign data for 50+ virtual devices",
        "🤖 AI-Powered: Tri-model architecture for comprehensive analysis",
        "⚡ High Performance: Built on FastAPI for low-latency inference",
        "🛡️ Secure: Anonymized patient data and modular design"
    ],
    "What We Monitor",
    [
        "❤️ Heart Rate",
        "🫁 Respiratory Rate",
        "🌡️ Temperature",
        "📊 Blood Oxygen (SpO2)",
        "🚨 Real-Time Alerts"
    ]
)

# Slide 4: AI Architecture - Tri-Model System
add_content_slide(prs, "AI/ML Model Architecture", [
    "• Edge Model (Isolation Forest): Instant anomaly detection on vital signs",
    "• TFT Model (Temporal Fusion Transformer): Predicts future health trends",
    "• CNN Model (ResNet50): Medical image analysis (X-ray anomaly detection)",
    "• Consolidated Decision: Aggregates model outputs for final health status",
    "• Result: Robust, multi-perspective health monitoring"
])

# Slide 5: Technical Architecture
add_two_column_slide(prs, "System Architecture",
    "Data Flow",
    [
        "1. Simulation generates real-time vital signs",
        "2. Data flows through Edge layer",
        "3. TFT service forecasts trends",
        "4. CNN processes medical images",
        "5. Results aggregated for final decision"
    ],
    "Layer Breakdown",
    [
        "Data Layer: IoT sensor streams",
        "Service Layer: AI orchestration",
        "Inference Engine: Model aggregation",
        "Presentation: Web dashboard",
        "Result: Normal/Anomaly Status"
    ]
)

# Slide 6: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    "Backend",
    [
        "Python 3.10+",
        "FastAPI (REST API)",
        "Uvicorn (ASGI server)",
        "PyTorch (AI/ML)",
        "Scikit-learn (ML algorithms)"
    ],
    "Frontend & Data",
    [
        "HTML5 & CSS3",
        "Vanilla JavaScript",
        "Real-time polling (3s intervals)",
        "CSV datasets",
        "Joblib & Pickle (serialization)"
    ]
)

# Slide 7: Data Flow & Backend Logic
add_content_slide(prs, "Data Flow & Backend Logic", [
    "• Simulator generates 50 virtual devices with realistic vital sign patterns",
    "• FastAPI orchestrates the simulation lifecycle and exposes REST endpoints",
    "• Frontend polls /devices endpoint every 3 seconds for updates",
    "• Each request passes vitals through Edge, TFT, and CNN services",
    "• Backend aggregates model outputs into Final Health Decision (Normal/Anomaly)"
])

# Slide 8: Dashboard & User Interface
add_content_slide(prs, "Dashboard & Visualization", [
    "• Modern, responsive monitoring interface",
    "• Displays real-time status for 20 active devices simultaneously",
    "• Color-coded health indicators: Normal (Green) vs. Anomaly (Red)",
    "• Live vital sign values for Heart Rate, SpO2, Respiratory Rate",
    "• Instant notifications for critical health alerts",
    "• 3-second polling refresh for up-to-date information"
])

# Slide 9: Project Structure & Setup
add_two_column_slide(prs, "Project Structure & Setup",
    "Directory Layout",
    [
        "/backend: FastAPI logic",
        "/models: Pre-trained weights",
        "/data: Healthcare datasets",
        "/frontend: Dashboard UI",
        "requirements.txt",
        "README.md & report.md"
    ],
    "Quick Start",
    [
        "1. Clone repository",
        "2. Create virtual environment",
        "3. pip install -r requirements.txt",
        "4. python backend/main.py",
        "5. Open dashboard in browser",
        "✓ System running!"
    ]
)

# Slide 10: Security, Benefits & Future
add_content_slide(prs, "Security, Benefits & Future Roadmap", [
    "🛡️ Security: Data privacy with anonymization, modular design for easy updates",
    "✅ Benefits: Real-time monitoring, multi-model redundancy, scalable architecture",
    "🚀 Future: Integrate real IoT devices, deploy edge computing, add predictive alerts",
    "💡 Innovation: Continuous model improvement, additional vital sign parameters",
    "🏥 Impact: Enhanced patient safety and healthcare provider decision support"
])

# Save presentation
output_path = r"C:\Users\subha\Downloads\iomt-ai-system\VitalGuard_AI_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
