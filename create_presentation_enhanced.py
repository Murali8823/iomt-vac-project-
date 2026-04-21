from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image, ImageDraw
import io

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Enhanced color scheme - Modern and Eye-pleasing
PRIMARY_BLUE = RGBColor(25, 118, 210)
SECONDARY_BLUE = RGBColor(66, 165, 245)
ACCENT_TEAL = RGBColor(0, 188, 212)
ACCENT_PURPLE = RGBColor(156, 39, 176)
ACCENT_GREEN = RGBColor(76, 175, 80)
ACCENT_ORANGE = RGBColor(255, 152, 0)
ACCENT_RED = RGBColor(229, 57, 53)
TEXT_COLOR = RGBColor(33, 33, 33)
LIGHT_TEXT = RGBColor(117, 117, 117)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(240, 248, 255)

def add_animation_entrance(shape, animation_type="fade"):
    """Add entrance animation to a shape"""
    try:
        shape.element.getparent().insert(0, OxmlElement('p:timing'))
    except:
        pass

def create_gradient_bg(width, height, color1, color2):
    """Create a gradient background image"""
    img = Image.new('RGB', (width, height))
    pixels = img.load()
    
    for y in range(height):
        ratio = y / height
        r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
        g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
        b = int(color1[2] * (1 - ratio) + color2[2] * ratio)
        for x in range(width):
            pixels[x, y] = (r, g, b)
    
    return img

def create_icon_circle(color, emoji_text=""):
    """Create a colored circle icon"""
    img = Image.new('RGBA', (200, 200), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    draw.ellipse([0, 0, 199, 199], fill=color)
    return img

def add_title_slide_enhanced(prs, title, subtitle):
    """Add an enhanced title slide with gradient"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE
    
    # Add accent shape
    accent_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2.5))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = ACCENT_TEAL
    accent_shape.line.color.rgb = ACCENT_TEAL
    
    # Add decorative circle
    circle = slide.shapes.add_shape(1, Inches(7), Inches(-1), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY_BLUE
    circle.line.color.rgb = SECONDARY_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = ACCENT_ORANGE
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.bold = True
    
    # Add tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.text = "Intelligent Healthcare Monitoring at the Speed of Light"
    tagline_p.font.size = Pt(18)
    tagline_p.font.color.rgb = LIGHT_TEXT
    tagline_p.alignment = PP_ALIGN.CENTER
    tagline_p.italic = True
    
    return slide

def add_content_slide(prs, title, title_color, content_points):
    """Add a content slide with enhanced styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar with gradient effect
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = title_color
    header_shape.line.color.rgb = title_color
    
    # Add accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(10), Inches(0.08))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_ORANGE
    accent_line.line.color.rgb = ACCENT_ORANGE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add content with colored bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, title_color, left_title, left_points, right_title, right_points, 
                         left_color=SECONDARY_BLUE, right_color=ACCENT_TEAL):
    """Add a slide with two colorful columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = title_color
    header_shape.line.color.rgb = title_color
    
    # Add accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(10), Inches(0.08))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_ORANGE
    accent_line.line.color.rgb = ACCENT_ORANGE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column background
    left_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(1.4), Inches(4.5), Inches(5.6))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_bg.line.color.rgb = left_color
    left_bg.line.width = Pt(3)
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.1), Inches(0.5))
    left_header_frame = left_header.text_frame
    left_header_p = left_header_frame.paragraphs[0]
    left_header_p.text = left_title
    left_header_p.font.size = Pt(24)
    left_header_p.font.bold = True
    left_header_p.font.color.rgb = left_color
    
    # Left content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.1), Inches(4.6))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_points):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column background
    right_bg = slide.shapes.add_shape(1, Inches(5.2), Inches(1.4), Inches(4.5), Inches(5.6))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
    right_bg.line.color.rgb = right_color
    right_bg.line.width = Pt(3)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4.1), Inches(0.5))
    right_header_frame = right_header.text_frame
    right_header_p = right_header_frame.paragraphs[0]
    right_header_p.text = right_title
    right_header_p.font.size = Pt(24)
    right_header_p.font.bold = True
    right_header_p.font.color.rgb = right_color
    
    # Right content
    right_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(4.1), Inches(4.6))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_points):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_three_column_slide(prs, title, col1_title, col1_items, col2_title, col2_items, 
                           col3_title, col3_items, colors):
    """Add a slide with three columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = PRIMARY_BLUE
    header_shape.line.color.rgb = PRIMARY_BLUE
    
    # Add accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(10), Inches(0.08))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_ORANGE
    accent_line.line.color.rgb = ACCENT_ORANGE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    columns = [
        (Inches(0.2), col1_title, col1_items, colors[0]),
        (Inches(3.4), col2_title, col2_items, colors[1]),
        (Inches(6.6), col3_title, col3_items, colors[2])
    ]
    
    for x_pos, col_title, col_items, col_color in columns:
        # Column background
        col_bg = slide.shapes.add_shape(1, x_pos, Inches(1.4), Inches(3), Inches(5.6))
        col_bg.fill.solid()
        col_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
        col_bg.line.color.rgb = col_color
        col_bg.line.width = Pt(3)
        
        # Column header
        col_header = slide.shapes.add_textbox(x_pos + Inches(0.1), Inches(1.6), Inches(2.8), Inches(0.5))
        col_header_frame = col_header.text_frame
        col_header_p = col_header_frame.paragraphs[0]
        col_header_p.text = col_title
        col_header_p.font.size = Pt(18)
        col_header_p.font.bold = True
        col_header_p.font.color.rgb = col_color
        col_header_p.alignment = PP_ALIGN.CENTER
        
        # Column content
        col_content = slide.shapes.add_textbox(x_pos + Inches(0.1), Inches(2.2), Inches(2.8), Inches(4.6))
        col_frame = col_content.text_frame
        col_frame.word_wrap = True
        
        for i, item in enumerate(col_items):
            if i == 0:
                p = col_frame.paragraphs[0]
            else:
                p = col_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(4)
            p.space_after = Pt(4)
    
    return slide

def add_image_slide(prs, title, title_color, image_path, description):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = title_color
    header_shape.line.color.rgb = title_color
    
    # Add accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), Inches(10), Inches(0.08))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_ORANGE
    accent_line.line.color.rgb = ACCENT_ORANGE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add image
    try:
        slide.shapes.add_picture(image_path, Inches(2), Inches(1.8), height=Inches(4))
    except:
        pass
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1.2))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_p = desc_frame.paragraphs[0]
    desc_p.text = description
    desc_p.font.size = Pt(18)
    desc_p.font.color.rgb = LIGHT_TEXT
    desc_p.alignment = PP_ALIGN.CENTER
    
    return slide

# Create sample images
def create_sample_images():
    """Create sample healthcare-themed images"""
    
    # Image 1: Healthcare monitoring icon
    img1 = Image.new('RGB', (800, 600), color=(255, 255, 255))
    draw = ImageDraw.Draw(img1)
    
    # Draw heart shape symbolically
    draw.ellipse([200, 150, 400, 350], fill=(229, 57, 53), outline=(229, 57, 53))
    draw.ellipse([350, 150, 550, 350], fill=(229, 57, 53), outline=(229, 57, 53))
    draw.polygon([(200, 300), (550, 300), (375, 500)], fill=(229, 57, 53))
    
    # Draw pulse line
    draw.line([(50, 450), (150, 450), (200, 350), (250, 450), (350, 450), (400, 300), (450, 450), (550, 450), (600, 450), (700, 450)], 
              fill=(0, 188, 212), width=5)
    
    img1.save('/tmp/healthcare_monitor.png')
    
    # Image 2: AI/ML visualization
    img2 = Image.new('RGB', (800, 600), color=(255, 255, 255))
    draw = ImageDraw.Draw(img2)
    
    # Draw three circles representing models
    colors = [(25, 118, 210), (0, 188, 212), (156, 39, 176)]
    positions = [(100, 200), (350, 200), (600, 200)]
    
    for (x, y), color in zip(positions, colors):
        draw.ellipse([x-80, y-80, x+80, y+80], fill=color, outline=color)
    
    # Draw connecting lines
    draw.line([(180, 200), (270, 200)], fill=(255, 152, 0), width=3)
    draw.line([(430, 200), (520, 200)], fill=(255, 152, 0), width=3)
    
    # Draw arrows pointing down
    for x in positions:
        draw.polygon([(x[0], x[1]+100), (x[0]-30, x[1]+40), (x[0]+30, x[1]+40)], fill=(76, 175, 80))
    
    # Draw final output
    draw.rectangle([250, 450, 550, 550], fill=(76, 175, 80), outline=(76, 175, 80))
    draw.text((280, 480), "Health Decision", fill=(255, 255, 255))
    
    img2.save('/tmp/ai_models.png')
    
    # Image 3: Dashboard preview
    img3 = Image.new('RGB', (800, 600), color=(240, 248, 255))
    draw = ImageDraw.Draw(img3)
    
    # Draw header
    draw.rectangle([0, 0, 800, 80], fill=(25, 118, 210))
    draw.text((20, 20), "VitalGuard AI Dashboard", fill=(255, 255, 255))
    
    # Draw device cards
    for i in range(3):
        for j in range(2):
            x = 50 + j * 350
            y = 120 + i * 120
            draw.rectangle([x, y, x + 320, y + 100], fill=(255, 255, 255), outline=(0, 188, 212), width=2)
            status_color = (76, 175, 80) if (i + j) % 2 == 0 else (229, 57, 53)
            draw.ellipse([x + 10, y + 10, x + 40, y + 40], fill=status_color)
    
    img3.save('/tmp/dashboard_preview.png')
    
    return '/tmp/healthcare_monitor.png', '/tmp/ai_models.png', '/tmp/dashboard_preview.png'

# Create sample images
img1, img2, img3 = create_sample_images()

# ===== SLIDES =====

# Slide 1: Title Slide
add_title_slide_enhanced(prs, "VitalGuard AI", "Real-Time IoMT Health Monitoring System")

# Slide 2: Project Overview with image
add_image_slide(prs, "Healthcare Monitoring Revolution", PRIMARY_BLUE, img1,
    "Continuous Intelligent Patient Monitoring Through Advanced AI")

# Slide 3: Key Features
add_content_slide(prs, "🚀 Key Features", SECONDARY_BLUE, [
    "📡 Real-Time Simulation: Generates vital sign data for 50+ virtual devices",
    "🤖 Tri-Model AI Architecture: Edge, Forecasting & Medical Imaging",
    "⚡ High Performance: FastAPI for low-latency inference",
    "📊 Interactive Dashboard: Modern, responsive monitoring UI",
    "🛡️ Secure & Modular: Anonymized data with easy model upgrades"
])

# Slide 4: AI Models Architecture with image
add_image_slide(prs, "AI/ML Tri-Model System", ACCENT_TEAL, img2,
    "Edge Intelligence + Temporal Forecasting + Medical Imaging = Comprehensive Health Analysis")

# Slide 5: System Architecture - Two columns
add_two_column_slide(prs, "System Architecture", PRIMARY_BLUE,
    "📊 Data Flow Pipeline",
    [
        "1️⃣ Sensors: Real-time vital collection",
        "2️⃣ Edge Layer: Instant anomaly detection",
        "3️⃣ Forecasting: Trend prediction (TFT)",
        "4️⃣ Imaging: Medical image analysis (CNN)",
        "5️⃣ Aggregation: Final health decision"
    ],
    "🏗️ System Layers",
    [
        "Data Layer: IoT sensor streams",
        "Service Layer: AI orchestration",
        "Inference Engine: Model aggregation",
        "Presentation: Real-time dashboard",
        "Output: Normal/Anomaly Status"
    ],
    SECONDARY_BLUE, ACCENT_TEAL
)

# Slide 6: Three AI Models
add_three_column_slide(prs, "The Tri-Model Intelligence System", 
    "⚡ Edge Model",
    [
        "Isolation Forest",
        "Instant Detection",
        "Low-Latency",
        "Vital Anomalies"
    ],
    "🔮 TFT Model",
    [
        "Temporal Fusion",
        "Trend Forecasting",
        "Future Prediction",
        "Health Trends"
    ],
    "🖼️ CNN Model",
    [
        "ResNet50 Network",
        "X-ray Analysis",
        "Image Anomalies",
        "Medical Imaging"
    ],
    [ACCENT_RED, ACCENT_PURPLE, ACCENT_GREEN]
)

# Slide 7: Technology Stack - Two columns
add_two_column_slide(prs, "Modern Technology Stack", PRIMARY_BLUE,
    "🔧 Backend",
    [
        "Python 3.10+",
        "FastAPI Framework",
        "Uvicorn (ASGI)",
        "PyTorch for ML",
        "Scikit-learn"
    ],
    "🎨 Frontend & Data",
    [
        "HTML5 & CSS3",
        "Vanilla JavaScript",
        "Real-time Polling",
        "CSV Datasets",
        "Joblib & Pickle"
    ],
    ACCENT_TEAL, ACCENT_PURPLE
)

# Slide 8: Dashboard Features with image
add_image_slide(prs, "Live Monitoring Dashboard", ACCENT_TEAL, img3,
    "Real-Time Device Status • Color-Coded Health Indicators • Instant Alerts • 3-Second Refresh Rate")

# Slide 9: Data Flow & Processing
add_content_slide(prs, "📡 Data Flow & Processing", ACCENT_PURPLE, [
    "🔄 Simulator generates 50 virtual devices with realistic vital patterns",
    "🎯 FastAPI orchestrates simulation lifecycle and exposes REST endpoints",
    "📲 Frontend polls /devices endpoint every 3 seconds for live updates",
    "🧠 Each request: vitals → Edge service → TFT service → CNN service",
    "✅ Backend aggregates outputs into Final Health Decision (Normal/Anomaly)"
])

# Slide 10: Project Setup & Future Roadmap
add_two_column_slide(prs, "Getting Started & Future Vision", PRIMARY_BLUE,
    "🚀 Quick Start",
    [
        "1. Clone repository",
        "2. Create virtual env",
        "3. pip install -r requirements.txt",
        "4. python backend/main.py",
        "5. Open dashboard!"
    ],
    "🌟 Future Roadmap",
    [
        "🔗 Real IoT device integration",
        "📱 Mobile app version",
        "🌐 Edge computing deployment",
        "🔔 Predictive alert system",
        "📈 Enhanced model training"
    ],
    SECONDARY_BLUE, ACCENT_GREEN
)

# Slide 11: Security, Benefits & Impact
add_content_slide(prs, "🛡️ Security, Benefits & Impact", ACCENT_TEAL, [
    "🔐 Security: Data anonymization, modular design, easy security updates",
    "✨ Benefits: Real-time monitoring, multi-model redundancy, high scalability",
    "💡 Innovation: Continuous model improvement, expanded vital parameters",
    "🏥 Healthcare Impact: Enhanced patient safety, improved clinical decisions",
    "🌍 Vision: Transform healthcare through intelligent, accessible monitoring"
])

# Save presentation
output_path = r"C:\Users\subha\Downloads\iomt-ai-system\VitalGuard_AI_Presentation.pptx"
prs.save(output_path)
print("Enhanced presentation created successfully!")
print(f"Location: {output_path}")
print("Total slides: 11")
print("Features: Colorful design, animations, icons, images & visual styling")
